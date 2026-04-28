"""
Update the Solution Challenge 2026 PPT template with project content.
Adds content text boxes below existing headings without changing the template format.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC = "[EXT] Solution Challenge 2026 - Prototype PPT Template.pptx"
DST = "[EXT] Solution Challenge 2026 - Prototype PPT Template.pptx"  # overwrite

FONT = "Google Sans"
BODY_COLOR = RGBColor(0x33, 0x33, 0x33)
ACCENT_COLOR = RGBColor(0x1A, 0x73, 0xE8)  # Google blue


def add_textbox(slide, left_in, top_in, width_in, height_in):
    return slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )


def set_para(para, text, size=11, bold=False, color=BODY_COLOR, bullet=False, indent=0):
    para.text = ""  # clear
    run = para.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if bullet:
        # add bullet using XML
        from pptx.oxml.ns import qn
        pPr = para._pPr if para._pPr is not None else para._p.get_or_add_pPr()
        # remove existing buChar/buNone
        for tag in ("buChar", "buNone", "buAutoNum"):
            for el in pPr.findall(qn(f"a:{tag}")):
                pPr.remove(el)
        from lxml import etree
        bu = etree.SubElement(pPr, qn("a:buChar"))
        bu.set("char", "•")
        pPr.set("indent", str(-Pt(0.18).emu) if False else "-228600")
        pPr.set("marL", str(228600 * (indent + 1)))
    return run


def write_lines(tf, lines, *, title=None, title_size=14, body_size=11):
    """lines: list of (text, kind) where kind in {'h','b','sub','plain'}"""
    tf.word_wrap = True
    tf.margin_left = Emu(91440)
    tf.margin_right = Emu(91440)
    tf.margin_top = Emu(45720)
    tf.margin_bottom = Emu(45720)

    first = True
    if title:
        p = tf.paragraphs[0]
        set_para(p, title, size=title_size, bold=True, color=ACCENT_COLOR)
        first = False

    for text, kind in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if kind == "h":
            set_para(p, text, size=body_size + 1, bold=True, color=ACCENT_COLOR)
        elif kind == "b":
            set_para(p, "•  " + text, size=body_size, color=BODY_COLOR)
        elif kind == "sub":
            set_para(p, "    – " + text, size=body_size - 1, color=BODY_COLOR)
        else:
            set_para(p, text, size=body_size, color=BODY_COLOR)


def main():
    prs = Presentation(SRC)

    # ==================== SLIDE 2: Team Details ====================
    slide = prs.slides[1]
    # find the team-details textbox and append values
    for shape in slide.shapes:
        if shape.has_text_frame and "Team name" in shape.text_frame.text:
            tf = shape.text_frame
            for para in tf.paragraphs:
                txt = para.text.strip()
                if txt.startswith("Team name:"):
                    # Append placeholder for user
                    if para.runs:
                        para.runs[0].text = "Team name:  Howlers"
                elif txt.startswith("Team leader name:"):
                    if para.runs:
                        para.runs[0].text = "Team leader name:  [Your Name]"
                elif txt.startswith("Problem Statement:"):
                    if para.runs:
                        para.runs[0].text = (
                            "Problem Statement:  Local NGOs and social groups collect "
                            "vital community-need data via paper surveys and field reports, "
                            "but the data stays scattered. Build a system to consolidate it, "
                            "surface the most urgent local needs, and intelligently match "
                            "available volunteers to the right tasks and areas."
                        )
            break

    # ==================== SLIDE 3: Brief about solution ====================
    slide = prs.slides[2]
    box = add_textbox(slide, 0.4, 1.6, 9.2, 3.8)
    write_lines(box.text_frame, [
        ("SamaajSetu  —  A Google-Cloud-native platform that turns scattered community-need data into intelligent volunteer action.", "plain"),
        ("", "plain"),
        ("What it does", "h"),
        ("Aggregates needs from paper surveys, mobile field reports, voice memos, SMS/IVR, and WhatsApp into a single, real-time view.", "b"),
        ("Uses Document AI to digitize handwritten/printed forms; Vertex AI Vector Search to match volunteers by skill + location + availability in milliseconds.", "b"),
        ("Dispatches the right volunteer to the right place at the right time via Firebase Cloud Messaging — works offline on entry-level Android phones.", "b"),
        ("Surfaces the most urgent local needs on a Looker Studio mission-control heatmap and predicts hotspots with BigQuery ML + Earth Engine.", "b"),
        ("Closes the loop with a 'Considerate Voice' feedback channel (Dialogflow CX + Gemini) that protects volunteer retention via the psychological contract.", "b"),
    ], body_size=11)

    # ==================== SLIDE 4: Opportunities (already has headings) ====================
    slide = prs.slides[3]
    box = add_textbox(slide, 0.4, 2.6, 9.2, 2.9)
    write_lines(box.text_frame, [
        ("How it differs from existing ideas", "h"),
        ("Most volunteer apps focus on hours-tracking; SamaajSetu focuses on need-aggregation + AI-driven matching, with paper surveys treated as a first-class data source.", "b"),
        ("Offline-first design works in low-bandwidth villages — most competitors assume always-online users.", "b"),
        ("How it solves the problem", "h"),
        ("Document AI converts paper-mountain into structured data; BigQuery centralizes it; Vertex AI Vector Search (KNN/ANN) finds the nearest skilled volunteer in <50 ms.", "b"),
        ("Pub/Sub + Cloud Workflows orchestrate dispatch end-to-end, replacing manual coordinator phone-trees during disasters.", "b"),
        ("USP", "h"),
        ("End-to-end on Google Cloud — single vendor, single bill, single IAM model — affordable for small NGOs via Google for Nonprofits credits.", "b"),
        ("Privacy-by-default: Cloud DLP, KMS-CMEK, and VPC-SC align with UN humanitarian data principles, including Community-Identifiable-Information protection.", "b"),
    ], body_size=10)

    # ==================== SLIDE 5: Features ====================
    slide = prs.slides[4]
    box = add_textbox(slide, 0.4, 1.7, 9.2, 3.7)
    write_lines(box.text_frame, [
        ("Paper-to-Data Digitization — Document AI Form Parser + Custom Extractor turns scanned surveys, handwritten notes, and PDFs into structured BigQuery rows.", "b"),
        ("Offline-First Field App — Flutter + Firestore offline persistence captures household visits, photos, GPS, and voice memos without internet; auto-syncs on reconnect.", "b"),
        ("AI-Powered Volunteer Matching — Vertex AI Embeddings + Vector Search rank volunteers by skill, geo-proximity, language, availability; Gemini re-ranks with explanations.", "b"),
        ("Real-Time Dispatch & Navigation — FCM push, SMS fallback, Routes API turn-by-turn directions, geo-verified check-ins.", "b"),
        ("Mission Control Dashboard — Looker Studio heatmap of urgent needs, KPIs, live activity feed for NGO coordinators.", "b"),
        ("Considerate Voice Channel — Dialogflow CX + Natural Language API for volunteer feedback; Gemini summarizes weekly themes for managers.", "b"),
        ("Multi-Channel Beneficiary Intake — SMS / IVR / WhatsApp / web forms in regional languages via Translation API + Speech-to-Text.", "b"),
        ("Crisis Mode — Cloud Workflows orchestrates disaster response with Earth Engine geospatial overlays and FCM topic broadcasts.", "b"),
        ("Predictive Insights — BigQuery ML forecasts demand and volunteer churn; pre-positions resources before incidents peak.", "b"),
        ("Privacy & Compliance — Cloud DLP tokenization, KMS-CMEK, VPC-SC perimeter, audit logs, role-based access.", "b"),
        ("Multi-Tenant SaaS — Identity Platform with phone-OTP, SAML/OIDC; per-NGO branding; AppSheet for no-code intake.", "b"),
        ("Donor Transparency Portal — Firebase Hosting public site with aggregated, redacted impact stats and live counters.", "b"),
    ], body_size=10)

    # ==================== SLIDE 6: Process flow / Use-case ====================
    slide = prs.slides[5]
    box = add_textbox(slide, 0.4, 1.65, 9.2, 3.8)
    write_lines(box.text_frame, [
        ("End-to-End Process Flow (full Mermaid diagrams in /docs/diagrams.md)", "h"),
        ("1. Capture — Field worker / paper / SMS / voice / web entry → Cloud Storage + Firestore (offline cache).", "b"),
        ("2. Digitize — Document AI (OCR + Form Parser), Speech-to-Text, Translation API normalize the input.", "b"),
        ("3. Privacy — Dataflow stream → Cloud DLP redacts/tokenizes PII; KMS-encrypted vault holds raw identifiers.", "b"),
        ("4. Store — Curated rows land in BigQuery (analytics) and Firestore (hot operational state).", "b"),
        ("5. Triage — Coordinator reviews in Triage Inbox (Kanban); urgency + program assigned; duplicates auto-detected.", "b"),
        ("6. Match — Vertex Embeddings → Vector Search (ANN top-K) → Gemini re-rank with 'why-matched' explanations.", "b"),
        ("7. Dispatch — FCM push + SMS fallback; Routes API navigation; geo-verified check-in via Firestore geofencing.", "b"),
        ("8. Execute — Volunteer performs task; submits photo + note + beneficiary signature (offline-queued).", "b"),
        ("9. Verify & Recognize — Beneficiary feedback; Gemini-personalized recognition push; certificates auto-issued.", "b"),
        ("10. Learn — BigQuery ML scores retention risk, demand forecast; Earth Engine flags emerging hotspots.", "b"),
        ("Key actors: Volunteer · Field Worker · Coordinator · NGO Admin · Beneficiary · Donor · Auditor · Govt Cell · Super Admin", "plain"),
    ], body_size=10)

    # ==================== SLIDE 7: Wireframes ====================
    slide = prs.slides[6]
    box = add_textbox(slide, 0.4, 1.7, 9.2, 3.7)
    write_lines(box.text_frame, [
        ("Volunteer Mobile App (Flutter, Android-first)", "h"),
        ("5-tab IA: Home (Today) · Discover (map) · My Tasks · Voice · Profile.", "b"),
        ("Key components: NeedCard (1-tap accept), WhyMatchedSheet, OfflineBadge, CheckInButton, RecognitionToast, AccessibilitySwitcher.", "b"),
        ("NGO Coordinator Web Console (Next.js)", "h"),
        ("Mission Control (heatmap + KPIs) · Triage Inbox (Kanban) · Matcher Studio (ranked list with explanation chips) · Programs · Volunteers · Surveys · Reports · Voice Channel.", "b"),
        ("Field Worker App", "h"),
        ("Form-first home, camera + voice + GPS as primary inputs, sync banner, beneficiary consent capture (signature + audio).", "b"),
        ("Beneficiary Touchpoints", "h"),
        ("SMS / IVR (Dialogflow CX), WhatsApp Business, lightweight web (<50 KB), public transparency portal.", "b"),
        ("Admin / Auditor / Donor / Govt Consoles", "h"),
        ("RBAC-scoped views: org settings & policies (Admin), audit-log explorer (Auditor), live counters & impact reports (Donor), cross-NGO situational view (Govt).", "b"),
        ("Design system 'Setu UI' shared across web + mobile (atomic components, design tokens, WCAG 2.2 AA, multilingual).", "plain"),
    ], body_size=10)

    # ==================== SLIDE 8: Architecture ====================
    slide = prs.slides[7]
    box = add_textbox(slide, 0.4, 1.7, 9.2, 3.7)
    write_lines(box.text_frame, [
        ("4-Layer GCP-Native Architecture", "h"),
        ("Ingest Layer — Document AI · Speech-to-Text · Firebase/Firestore · AppSheet · Cloud Storage · Pub/Sub · SMS Gateway (Cloud Run).", "b"),
        ("Process Layer — Dataflow streaming · Cloud DLP · Translation API · Natural Language API · Geocoding · Cloud Workflows · Eventarc.", "b"),
        ("Intelligence Layer — Vertex AI Embeddings · Vertex AI Vector Search (ScaNN, ANN <50 ms) · Vertex AI Gemini · BigQuery ML · Earth Engine · Dialogflow CX.", "b"),
        ("Engage Layer — Firebase Cloud Messaging · Maps Platform · Routes API · Looker Studio · Firebase Hosting · Identity Platform · Apigee.", "b"),
        ("11 Domain Microservices on Cloud Run (DDD bounded contexts)", "h"),
        ("identity-svc · intake-svc · digitization-svc · need-svc · volunteer-svc · matching-svc · dispatch-svc · engagement-svc · analytics-svc · trust-svc · billing-svc.", "b"),
        ("Cross-Cutting", "h"),
        ("Multi-tenant via org_id partitioning · RBAC + Identity-Aware Proxy · KMS-CMEK on every datastore · VPC Service Controls perimeter · Cloud Audit Logs.", "b"),
        ("Event-driven via Pub/Sub topics (needs.* volunteers.* match.* alerts.* audit.*); orchestration via Cloud Workflows; CI/CD via Cloud Build → Cloud Deploy.", "b"),
        ("Full architecture & deployment diagrams: /docs/product_architecture.md and /docs/diagrams.md.", "plain"),
    ], body_size=10)

    # ==================== SLIDE 9: Technologies ====================
    slide = prs.slides[8]
    box = add_textbox(slide, 0.4, 2.3, 9.2, 3.2)
    write_lines(box.text_frame, [
        ("Google AI / ML", "h"),
        ("Vertex AI (Gemini, Embeddings, Vector Search, Pipelines, Experiments) · Document AI · Speech-to-Text · Translation API · Natural Language API · Dialogflow CX · BigQuery ML · Earth Engine · Firebase ML.", "b"),
        ("Compute & Eventing", "h"),
        ("Cloud Run · Cloud Functions · Cloud Workflows · Cloud Tasks · Pub/Sub · Eventarc · Apigee · API Gateway · Cloud Build · Cloud Deploy · Artifact Registry.", "b"),
        ("Data & Storage", "h"),
        ("BigQuery · Firestore · Cloud Storage · Cloud SQL / AlloyDB · Data Catalog · Dataflow · Looker Studio · Looker.", "b"),
        ("Engagement & Maps", "h"),
        ("Firebase (Auth, Firestore, Hosting, Remote Config, Crashlytics) · Firebase Cloud Messaging · Maps Platform · Routes API · Geocoding API.", "b"),
        ("Security & Compliance", "h"),
        ("Identity Platform · Identity-Aware Proxy · Cloud DLP · Cloud KMS (CMEK) · VPC Service Controls · Assured Workloads · Secret Manager · Cloud Audit Logs · Access Transparency.", "b"),
        ("Frontend", "h"),
        ("Flutter (mobile) · Next.js + React + Tailwind + shadcn/ui (web) · Riverpod · go_router · Material 3 · WCAG 2.2 AA accessibility.", "b"),
        ("Observability", "h"),
        ("Cloud Logging · Cloud Trace · Cloud Monitoring · Error Reporting · SLO dashboards.", "b"),
    ], body_size=10)

    # ==================== SLIDE 10: Estimated cost ====================
    slide = prs.slides[9]
    box = add_textbox(slide, 0.4, 1.7, 9.2, 3.8)
    write_lines(box.text_frame, [
        ("MVP Pilot (1 NGO, ~500 volunteers, ~2,000 needs/month)", "h"),
        ("Document AI (Form Parser): ~10,000 pages × $0.065  ≈ $650/mo", "b"),
        ("Vertex AI Vector Search: small index (10K vectors) + queries  ≈ $80/mo", "b"),
        ("Vertex AI Gemini (re-rank + summaries): light usage  ≈ $40/mo", "b"),
        ("BigQuery: <100 GB storage + queries (free tier covers most)  ≈ $20/mo", "b"),
        ("Firestore + Firebase Auth + FCM: within Spark/Blaze low-tier  ≈ $30/mo", "b"),
        ("Cloud Run (11 svc, scale-to-zero) + Pub/Sub + Workflows  ≈ $60/mo", "b"),
        ("Maps Platform (Geocoding + Routes): ~50K calls  ≈ $100/mo", "b"),
        ("Cloud DLP + KMS + Audit Logs + Monitoring  ≈ $40/mo", "b"),
        ("Estimated MVP total:  ~$1,000 / month  (further reduced via Google for Nonprofits credits — typically $2,000–$10,000 annual credit)", "h"),
        ("", "plain"),
        ("Scale Estimate (10 NGOs, 10,000 volunteers, 50,000 needs/month)  ≈ $4,000–$6,000 / month", "b"),
        ("Cost-saving levers: scale-to-zero Cloud Run, BigQuery committed slots, sustained-use discounts, AppSheet free tier for small partners, Google for Nonprofits.", "plain"),
    ], body_size=10)

    # ==================== SLIDE 11: MVP Snapshots (placeholder note) ====================
    slide = prs.slides[10]
    box = add_textbox(slide, 0.4, 1.7, 9.2, 3.7)
    write_lines(box.text_frame, [
        ("MVP Surfaces (insert screenshots here once built)", "h"),
        ("1.  Volunteer Mobile App — Home tab with 3 personalized NeedCards + WhyMatched explanation chips.", "b"),
        ("2.  Field Worker Survey Form — offline capture screen showing photo + GPS + sync banner.", "b"),
        ("3.  Coordinator Mission Control — Looker Studio heatmap + KPI tiles (Open Needs, Active Volunteers, SLA Breaches).", "b"),
        ("4.  Triage Inbox — Kanban board with Submitted → Triaged → Published → Matched → In Progress columns.", "b"),
        ("5.  Matcher Studio — ranked volunteer list with skill-match, distance, language chips and dispatch button.", "b"),
        ("6.  Document AI demo — scanned paper survey → structured JSON extraction in BigQuery.", "b"),
        ("7.  Voice Channel — Dialogflow CX chat for volunteer concerns + Gemini-summarized weekly digest.", "b"),
        ("8.  Crisis Mode — district heatmap with Earth Engine flood overlay and bulk dispatch panel.", "b"),
        ("9.  Donor Transparency Portal — public Firebase-hosted page with live impact counters.", "b"),
    ], body_size=10)

    # ==================== SLIDE 12: Future Development ====================
    slide = prs.slides[11]
    box = add_textbox(slide, 0.4, 1.7, 9.2, 3.7)
    write_lines(box.text_frame, [
        ("Near-term (3–6 months)", "h"),
        ("Volunteer-skill verification via partner certifications (e.g., Red Cross First Aid) — auto-imported badges.", "b"),
        ("Federated learning for on-device need-classification, reducing data egress for privacy-sensitive contexts.", "b"),
        ("WhatsApp Business deep-integration for two-way beneficiary follow-ups in rich media.", "b"),
        ("Mid-term (6–12 months)", "h"),
        ("Cross-NGO volunteer marketplace with consent-based skill-sharing during regional crises.", "b"),
        ("Predictive deployment using Earth Engine + weather APIs to pre-position before disasters strike.", "b"),
        ("Voice-first interface for low-literacy users via Vertex AI conversational agents in 20+ Indian languages.", "b"),
        ("Long-term (12+ months)", "h"),
        ("Block-chain-anchored impact ledger (signed by Cloud KMS) for verifiable donor reporting.", "b"),
        ("Agentic-AI 'Program Co-Pilot' that proposes program adjustments by analyzing outcome data and similar programs across the network.", "b"),
        ("Open API + AppSheet templates marketplace so any NGO can extend the platform without engineers.", "b"),
        ("Research collaborations to publish anonymized governance-effectiveness datasets (UN-aligned) — addresses the 'paucity of comparative research' gap.", "b"),
    ], body_size=10)

    # Slide 13 (links) — leave for user to fill with GitHub / Demo / MVP / Prototype URLs
    # Slides 14, 15 are picture-only ending slides — leave untouched.

    prs.save(DST)
    print(f"Saved updated PPT to: {DST}")


if __name__ == "__main__":
    main()
